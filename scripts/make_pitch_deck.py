"""Generate the CutPilot hackathon pitch deck.

Colors + type match the review UI at ui/index.html (Inter + JetBrains Mono,
#f5f5f7 page, white cards with #eeeeee borders, #fc3f1d brand red, black primary).

Run: `python scripts/make_pitch_deck.py`
Writes: `deliverables/cutpilot_pitch.pptx` (overwrites).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = REPO_ROOT / "deliverables"
OUT_PATH = OUT_DIR / "cutpilot_pitch.pptx"

# Palette lifted from ui/index.html
PAGE = RGBColor(0xF5, 0xF5, 0xF7)           # body background
CARD = RGBColor(0xFF, 0xFF, 0xFF)           # .card
CARD_BORDER = RGBColor(0xEE, 0xEE, 0xEE)    # .card border
HAIRLINE = RGBColor(0xF0, 0xF0, 0xF0)       # dividers, score-bar track
INK_900 = RGBColor(0x00, 0x00, 0x00)        # primary text, btn-primary
INK_800 = RGBColor(0x1A, 0x1A, 0x1A)
INK_700 = RGBColor(0x2B, 0x2B, 0x2B)
INK_500 = RGBColor(0x6B, 0x6B, 0x6B)        # muted
INK_400 = RGBColor(0x8A, 0x8A, 0x8A)
INK_300 = RGBColor(0xC4, 0xC4, 0xC4)
INK_200 = RGBColor(0xE3, 0xE3, 0xE3)
INK_100 = RGBColor(0xEE, 0xEE, 0xEE)
INK_50 = RGBColor(0xF5, 0xF5, 0xF7)
RED = RGBColor(0xFC, 0x3F, 0x1D)            # yred-500, brand-dot, progress-fill
RED_DARK = RGBColor(0xC4, 0x2D, 0x16)       # yred-700
YELLOW = RGBColor(0xFF, 0xCC, 0x00)         # yyellow-500

SANS = "Inter"   # UI primary; falls back to Helvetica Neue on systems without Inter
MONO = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def _textbox(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK_900, align=PP_ALIGN.LEFT, font=SANS):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _bullets(slide, left, top, width, height, items, *, size=18,
             color=INK_800, bullet_color=RED, line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        dot = p.add_run()
        dot.text = "●  "
        dot.font.name = SANS
        dot.font.size = Pt(size)
        dot.font.bold = True
        dot.font.color.rgb = bullet_color
        body = p.add_run()
        body.text = item
        body.font.name = SANS
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return tb


def _card(slide, left, top, width, height, *, fill=CARD, border=CARD_BORDER,
          border_pt=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # Tighten the corner radius on the rounded rect (default is chunky)
    try:
        s.adjustments[0] = 0.06
    except Exception:
        pass
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(border_pt)
    s.shadow.inherit = False
    return s


def _chip(slide, left, top, text, *, bg=INK_900, fg=CARD, size=11, pad=Inches(0.14)):
    # Approximate width — chips are short labels
    w = Inches(max(0.9, 0.11 * len(text) + 0.35))
    h = Inches(0.34)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    try:
        s.adjustments[0] = 0.5
    except Exception:
        pass
    s.fill.solid()
    s.fill.fore_color.rgb = bg
    s.line.fill.background()
    s.shadow.inherit = False
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = MONO
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return w


def _brand_dot(slide, left, top, size=Inches(0.22)):
    # Solid red disc with a soft red halo — matches .brand-dot in the UI.
    halo_off = Inches(0.06)
    halo = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left - halo_off, top - halo_off,
        size + halo_off * 2, size + halo_off * 2,
    )
    halo.fill.solid()
    halo.fill.fore_color.rgb = RGBColor(0xFD, 0xD9, 0xD1)  # rgba(252,63,29,0.12) on white
    halo.line.fill.background()
    halo.shadow.inherit = False
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = RED
    dot.line.fill.background()
    dot.shadow.inherit = False


def _title(slide, title, subtitle=None):
    # CutPilot brand lockup, top-left — mirrors the UI header
    _brand_dot(slide, Inches(0.6), Inches(0.52))
    _textbox(slide, Inches(0.95), Inches(0.4), Inches(6), Inches(0.5),
             "CutPilot", size=16, bold=True, color=INK_900)
    # Right-aligned track tag
    _textbox(slide, Inches(9.6), Inches(0.45), Inches(3.3), Inches(0.35),
             "Track A · Nemotron Developer Days",
             size=10, color=INK_500, align=PP_ALIGN.RIGHT, font=MONO)
    # Thin separator (matches the UI's subtle card dividers)
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.0), Inches(12.15), Inches(0.01),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = HAIRLINE
    sep.line.fill.background()
    sep.shadow.inherit = False
    # Page title + subtitle
    _textbox(slide, Inches(0.6), Inches(1.25), Inches(12), Inches(0.8),
             title, size=36, bold=True, color=INK_900)
    if subtitle:
        _textbox(slide, Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
                 subtitle, size=16, color=INK_500)


def _footer(slide, page, total):
    _textbox(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "cutpilot  ·  agentic long-video → 3 vertical highlights",
             size=10, color=INK_400)
    _textbox(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
             f"{page:02d} / {total:02d}",
             size=10, color=INK_400, align=PP_ALIGN.RIGHT, font=MONO)


def _label(slide, left, top, width, height, text, *, size=14, bold=False,
           color=INK_900, align=PP_ALIGN.CENTER, font=SANS):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


# ---------- slides ----------

def slide_title(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)

    # Brand lockup, a bit larger on the cover
    _brand_dot(s, Inches(0.9), Inches(0.92), size=Inches(0.36))
    _textbox(s, Inches(1.45), Inches(0.8), Inches(10), Inches(0.6),
             "CutPilot", size=22, bold=True, color=INK_900)

    _textbox(s, Inches(9.6), Inches(0.88), Inches(3.3), Inches(0.4),
             "Track A · Nemotron Developer Days",
             size=11, color=INK_500, align=PP_ALIGN.RIGHT, font=MONO)

    # Big headline, UI-style — three pre-broken lines so PPT doesn't re-wrap
    _textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.95),
             "An agentic editor that turns",
             size=44, bold=True, color=INK_500)
    _textbox(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.95),
             "long-form video into three",
             size=44, bold=True, color=INK_900)
    _textbox(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.95),
             "vertical highlights.",
             size=44, bold=True, color=INK_900)

    # Red underline accent, below the last headline line
    acc = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.38),
        Inches(1.4), Inches(0.08),
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = RED
    acc.line.fill.background()
    acc.shadow.inherit = False

    # Stack chips, like the UI's status pills
    left = Inches(0.9)
    top = Inches(4.55)
    chip_specs = [
        ("Whisper NIM", INK_900, CARD),
        ("Nemotron Nano 2 VL", INK_900, CARD),
        ("Nemotron Text", INK_900, CARD),
        ("NeMo Agent Toolkit", RED, CARD),
        ("NVIDIA Brev H100", INK_800, CARD),
    ]
    x = left
    for text, bg, fg in chip_specs:
        w = _chip(s, x, top, text, bg=bg, fg=fg, size=11)
        x += w + Inches(0.12)

    # Meta line
    _textbox(s, Inches(0.9), Inches(6.5), Inches(12), Inches(0.4),
             "Nemotron Developer Days Seoul · 2026",
             size=12, color=INK_500, font=MONO)
    _textbox(s, Inches(0.9), Inches(6.85), Inches(12), Inches(0.4),
             "sergey@lablup.com",
             size=12, color=INK_500, font=MONO)


def slide_problem(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)
    _title(s, "The problem",
           "Long-form video is where the ideas live. Short-form is where the audience is.")

    _card(s, Inches(0.6), Inches(2.7), Inches(12.15), Inches(3.2))
    _bullets(s, Inches(1.0), Inches(2.95), Inches(11.3), Inches(3.0), [
        "Creators spend 2–4 hours scrubbing a 60-minute video to pull three postable clips.",
        "Heuristic tools (loudness peaks, keyword hits) can’t explain their picks — editors don’t trust black boxes.",
        "Raw clips aren’t broadcast-ready: need 9:16 reframe, burned captions, tight boundaries that don’t clip mid-sentence.",
        "Teams want a reasoning trace: why this clip, what did you skip, at what confidence?",
    ], size=19)

    _card(s, Inches(0.6), Inches(6.1), Inches(12.15), Inches(0.75),
          fill=INK_900, border=INK_900)
    _textbox(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.6),
             "An editor that picks fast, explains itself, and ships the clip — not a dashboard of features.",
             size=16, bold=True, color=CARD)
    _footer(s, 2, total)


def slide_solution(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)
    _title(s, "The solution",
           "One command in. Three postable clips out. With reasoning.")

    # Black terminal-style card (like .btn-primary)
    _card(s, Inches(0.6), Inches(2.7), Inches(12.15), Inches(1.35),
          fill=INK_900, border=INK_900)
    _textbox(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.5),
             "$  cutpilot  https://youtu.be/<id>          # or a local .mp4",
             size=20, bold=True, color=YELLOW, font=MONO)
    _textbox(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(0.5),
             "$  cutpilot-serve                           # FastAPI + review UI at 127.0.0.1:8080",
             size=16, color=CARD, font=MONO)

    _card(s, Inches(0.6), Inches(4.3), Inches(12.15), Inches(2.6))
    _bullets(s, Inches(1.0), Inches(4.5), Inches(11.3), Inches(2.4), [
        "Input: a 5–90 min podcast, lecture, keynote, or interview — local file or URL.",
        "Output: three 30–60 s vertical clips (1080×1920) with burned-in captions + a stitched highlights.mp4.",
        "Each clip ships a JSON manifest: hook, multi-sentence rationale, 4-axis rubric (hook / self-contained / length / visual).",
        "Review UI mirrors this deck: white cards, black primary, red live-state dot, reasoning trace one click away.",
    ], size=17)

    _footer(s, 3, total)


def slide_demo_flow(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)
    _title(s, "Demo flow", "Four stages. Streamed to the UI as they complete.")

    stages = [
        ("1", "Ingest", "yt-dlp  ·  local"),
        ("2", "Transcribe", "Whisper NIM\nword-level"),
        ("3", "Scout", "Nemotron Nano 2 VL\n5–10 candidates + scores"),
        ("4", "Render", "ffmpeg:  cut  ·  9:16  ·  captions"),
    ]
    left = Inches(0.6)
    top = Inches(2.85)
    box_w = Inches(2.9)
    box_h = Inches(2.0)
    gap = Inches(0.23)
    for i, (num, head, sub) in enumerate(stages):
        x = left + (box_w + gap) * i
        _card(s, x, top, box_w, box_h)
        # Numbered brand-dot style label
        _brand_dot(s, x + Inches(0.25), top + Inches(0.25), size=Inches(0.22))
        _label(s, x + Inches(0.55), top + Inches(0.23), Inches(1.0), Inches(0.3),
               num, size=12, bold=True, color=INK_900, align=PP_ALIGN.LEFT, font=MONO)
        _label(s, x, top + Inches(0.75), box_w, Inches(0.45),
               head, size=20, bold=True, color=INK_900)
        _label(s, x + Inches(0.2), top + Inches(1.25), box_w - Inches(0.4), Inches(0.7),
               sub, size=13, color=INK_500)
        if i < len(stages) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w + Inches(0.025),
                top + box_h / 2 - Inches(0.08),
                Inches(0.18),
                Inches(0.16),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = INK_300
            arrow.line.fill.background()
            arrow.shadow.inherit = False

    _textbox(s, Inches(0.6), Inches(5.0), Inches(12.15), Inches(0.4),
             "Stage transitions surface live via an on_stage callback → FastAPI → UI polls GET /runs/{id}.",
             size=12, color=INK_500, font=MONO)

    _card(s, Inches(0.6), Inches(5.45), Inches(12.15), Inches(1.4))
    _textbox(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.4),
             "Deliverables per run", size=15, bold=True, color=INK_900)
    _bullets(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.0), [
        "outputs/<run>/clip_{1,2,3}.mp4                     —  vertical, captioned, broadcast-ready",
        "outputs/<run>/clip_{1,2,3}.manifest.json    —  hook, rationale, rubric scores, timestamps",
        "outputs/<run>/highlights.mp4                          —  stitched 90–180 s reel for one-tap posting",
    ], size=12, line_spacing=1.2, bullet_color=INK_900)
    _footer(s, 4, total)


def slide_architecture(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)
    _title(s, "Architecture",
           "Three layers. Top decides the flow, middle decides the content, bottom does the work.")

    # Three vertically-stacked bands. Each is one sentence of plain English +
    # a short example of what lives there. No API names, no code.
    bands = [
        ("1",
         "Pipeline",
         "Deterministic Python that moves a video through the stages.",
         "Download  →  transcribe  →  scout  →  edit  →  save"),
        ("2",
         "Agents",
         "Two LLMs make every content decision — what to clip and why.",
         "Scout picks the moments.      Editor writes the cut plan."),
        ("3",
         "Tools",
         "Small utilities the agents call to actually touch files.",
         "cut  ·  crop  ·  captions  ·  splice  ·  transcript"),
    ]
    top = Inches(2.75)
    band_h = Inches(1.25)
    gap = Inches(0.2)
    for i, (num, head, sub, example) in enumerate(bands):
        y = top + (band_h + gap) * i
        _card(s, Inches(0.6), y, Inches(12.15), band_h)

        # Number badge on the left
        _brand_dot(s, Inches(0.95), y + Inches(0.4), size=Inches(0.42))
        _label(s, Inches(0.95), y + Inches(0.4), Inches(0.42), Inches(0.42),
               num, size=20, bold=True, color=CARD, align=PP_ALIGN.CENTER, font=SANS)

        # Layer name
        _textbox(s, Inches(1.9), y + Inches(0.22), Inches(3.5), Inches(0.5),
                 head, size=22, bold=True, color=INK_900)
        # Plain-English subhead
        _textbox(s, Inches(1.9), y + Inches(0.63), Inches(10.5), Inches(0.4),
                 sub, size=14, color=INK_500)
        # Concrete example (mono, muted) — keeps the slide concrete without code
        _textbox(s, Inches(1.9), y + Inches(0.93), Inches(10.5), Inches(0.3),
                 example, size=12, color=INK_700, font=MONO)

    _footer(s, 5, total)


def slide_nemo_stack(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)
    _title(s, "NeMo stack usage",
           "Three NVIDIA NIMs. Each does one job. The toolkit glues them together.")

    # Three role-based cards — role first (Ears / Eyes / Brain), name second,
    # then one line of what it produces. No config params, no SDK details.
    cards = [
        ("Ears",  "Whisper NIM",
         "Listens to the audio and writes down every word with a timestamp."),
        ("Eyes",  "Nemotron Nano 2 VL",
         "Watches the video and describes what's visually happening in each window."),
        ("Brain", "Nemotron Text",
         "Reads the transcript + visual notes and decides which 3 clips to ship."),
    ]
    top = Inches(2.75)
    box_w = Inches(3.95)
    box_h = Inches(2.2)
    gap = Inches(0.15)
    for i, (role, name, body) in enumerate(cards):
        x = Inches(0.6) + (box_w + gap) * i
        _card(s, x, top, box_w, box_h)
        # Role label at top — the plain-English hook
        _label(s, x, top + Inches(0.22), box_w, Inches(0.45),
               role, size=13, bold=True, color=RED, font=MONO)
        # NIM name
        _label(s, x, top + Inches(0.6), box_w, Inches(0.5),
               name, size=20, bold=True, color=INK_900)
        # One-sentence description
        _textbox(s, x + Inches(0.25), top + Inches(1.2), box_w - Inches(0.5), Inches(1.0),
                 body, size=14, color=INK_800, align=PP_ALIGN.CENTER)

    # Toolkit contributions — three plain-English bullets, no API names
    _card(s, Inches(0.6), Inches(5.2), Inches(12.15), Inches(1.65))
    _textbox(s, Inches(0.85), Inches(5.32), Inches(11.6), Inches(0.45),
             "NeMo Agent Toolkit — what it gives us",
             size=16, bold=True, color=INK_900)
    _bullets(s, Inches(0.85), Inches(5.8), Inches(12), Inches(1.0), [
        "One YAML file describes the whole agent workflow — swap models without touching code.",
        "Tools register themselves — the toolkit discovers them automatically at install time.",
        "Ships with the agent types we need out of the box (tool-calling, sequential orchestration).",
    ], size=13, line_spacing=1.3, bullet_color=RED)

    _footer(s, 6, total)


def slide_innovation(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)
    _title(s, "What makes it different",
           "Four ideas most clippers skip — each one is easy to explain on stage.")

    # Four ideas, two columns × two rows. Each card: a 2-5 word headline a
    # presenter can say out loud, then one sentence of payoff. No jargon.
    items = [
        ("It explains every pick.",
         "Every clip ships with a rationale and a score. An editor can defend it to their boss."),
        ("It sees the video, not just hears it.",
         "A visual model scores what's on screen so we don't pick the most boring moment with the loudest audio."),
        ("It reads the whole transcript.",
         "A language model picks moments by meaning — not by volume peaks or keyword hits."),
        ("It ships ready-to-post clips.",
         "Vertical 9:16 crop and the stitched highlight reel come out of the same run. No second tool."),
    ]
    top = Inches(2.75)
    card_w = Inches(6.0)
    card_h = Inches(1.85)
    gap = Inches(0.15)
    for i, (head, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (card_w + gap) * col
        y = top + (card_h + gap) * row
        _card(s, x, y, card_w, card_h)
        # Red accent bar (matches brand language on the cover)
        acc = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.35), y + Inches(0.35), Inches(0.6), Inches(0.06),
        )
        acc.fill.solid()
        acc.fill.fore_color.rgb = RED
        acc.line.fill.background()
        acc.shadow.inherit = False
        _textbox(s, x + Inches(0.35), y + Inches(0.5), card_w - Inches(0.7), Inches(0.6),
                 head, size=22, bold=True, color=INK_900)
        _textbox(s, x + Inches(0.35), y + Inches(1.1), card_w - Inches(0.7), Inches(0.7),
                 body, size=14, color=INK_500)

    _footer(s, 7, total)


def slide_impact(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)
    _title(s, "Impact & scalability",
           "Real users. Real scaling story. Agentic ≠ opaque.")

    _card(s, Inches(0.6), Inches(2.75), Inches(6.0), Inches(2.6))
    _textbox(s, Inches(0.9), Inches(2.9), Inches(5.5), Inches(0.4),
             "Who this is for", size=16, bold=True, color=INK_900)
    _bullets(s, Inches(0.9), Inches(3.35), Inches(5.6), Inches(2.0), [
        "Creators / podcasters clipping long episodes.",
        "Conference organizers posting daily highlights.",
        "Educators turning 90-min lectures into promos.",
        "Marketing teams surfacing customer interviews at scale.",
    ], size=13, line_spacing=1.25)

    _card(s, Inches(6.75), Inches(2.75), Inches(6.0), Inches(2.6))
    _textbox(s, Inches(7.05), Inches(2.9), Inches(5.5), Inches(0.4),
             "How it scales", size=16, bold=True, color=INK_900)
    _bullets(s, Inches(7.05), Inches(3.35), Inches(5.6), Inches(2.0), [
        "Horizontal: add NIM replicas behind each endpoint — Scout is stateless per run.",
        "Model swap: one-line YAML edit; tools untouched.",
        "Self-hosted NIM on Brev today → same code to hosted build.nvidia.com NIMs.",
        "Air-gapped enterprise: no third-party editing SaaS in the loop.",
    ], size=13, line_spacing=1.25, bullet_color=INK_900)

    # Why-agentic panel
    _card(s, Inches(0.6), Inches(5.55), Inches(12.15), Inches(1.35))
    _textbox(s, Inches(0.9), Inches(5.65), Inches(11.6), Inches(0.4),
             "Why an agentic editor beats a heuristic one — one line per stakeholder",
             size=15, bold=True, color=INK_900)
    _bullets(s, Inches(0.9), Inches(6.05), Inches(12), Inches(0.9), [
        "Editor:   every clip ships with a rationale. You can defend the pick to your boss.",
        "Platform: the reasoning trace is an audit log. Regulated industries can keep it.",
        "Buyer:    config-driven — change the rubric (‘must mention the product’) in YAML, not in code.",
    ], size=12, line_spacing=1.25, bullet_color=RED)
    _footer(s, 8, total)


def slide_demo_evidence(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)
    _title(s, "Live demo evidence",
           "Real output from a 43-minute GTC DC 2025 Healthcare keynote.")

    _textbox(s, Inches(0.6), Inches(2.7), Inches(12.15), Inches(0.35),
             "outputs/e2e-sergey/  —  three clips produced end-to-end, verbatim from clip_*.manifest.json, no hand edits.",
             size=11, color=INK_500, font=MONO)

    clips = [
        ("Clip 1  ·  0:00 – 0:55",
         "What if AI could design life-saving drugs in minutes?",
         [("hook", 5), ("self-contained", 4), ("length", 5), ("visual", 4)],
         "Opening segment frames GTC as a transformative reunion, capturing attention with a bold vision of AI reshaping healthcare."),
        ("Clip 2  ·  5:00 – 6:00",
         "How can AI turn genomic data into real-time diagnoses for newborns?",
         [("hook", 5), ("self-contained", 4), ("length", 4), ("visual", 4)],
         "Captures breakthroughs in genomic sequencing — incl. a Guinness World Record for neonatal diagnosis. Highly postable."),
        ("Clip 3  ·  10:00 – 11:15",
         "What if proteins could be designed like software?",
         [("hook", 4), ("self-contained", 5), ("length", 5), ("visual", 3)],
         "Speaker introduces OpenFold3 and NIM microservices; 75 s window stays within the 20–90 s limit."),
    ]
    top = Inches(3.1)
    card_h = Inches(1.2)
    gap = Inches(0.12)
    for i, (label, hook, scores, rationale) in enumerate(clips):
        y = top + (card_h + gap) * i
        _card(s, Inches(0.6), y, Inches(12.15), card_h)
        _textbox(s, Inches(0.9), y + Inches(0.1), Inches(3.5), Inches(0.35),
                 label, size=11, bold=True, color=RED, font=MONO)
        _textbox(s, Inches(0.9), y + Inches(0.38), Inches(12), Inches(0.4),
                 hook, size=16, bold=True, color=INK_900)
        # score mini-bars (score / 5)
        bar_top = y + Inches(0.82)
        bar_x = Inches(0.9)
        for name, val in scores:
            _textbox(s, bar_x, bar_top - Inches(0.02), Inches(1.4), Inches(0.25),
                     name, size=9, color=INK_500, font=MONO)
            track = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                bar_x, bar_top + Inches(0.22), Inches(1.2), Inches(0.06),
            )
            track.fill.solid()
            track.fill.fore_color.rgb = HAIRLINE
            track.line.fill.background()
            track.shadow.inherit = False
            fill = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                bar_x, bar_top + Inches(0.22),
                Inches(1.2 * (val / 5.0)), Inches(0.06),
            )
            fill.fill.solid()
            fill.fill.fore_color.rgb = INK_900
            fill.line.fill.background()
            fill.shadow.inherit = False
            bar_x += Inches(1.55)
        _textbox(s, Inches(7.7), y + Inches(0.78), Inches(5.0), Inches(0.4),
                 rationale, size=10, color=INK_500)

    _footer(s, 9, total)


def slide_completeness(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)
    _title(s, "Completeness & quality",
           "Track A asks for a complete system. Here’s how it stacks up.")

    stats = [
        ("3,020", "lines of Python"),
        ("28", "source files"),
        ("19", "test files"),
        ("13 · 5 · 1", "unit · integration · e2e"),
    ]
    top = Inches(2.75)
    box_w = Inches(2.95)
    box_h = Inches(1.35)
    gap = Inches(0.18)
    for i, (big, small) in enumerate(stats):
        x = Inches(0.6) + (box_w + gap) * i
        _card(s, x, top, box_w, box_h)
        _label(s, x, top + Inches(0.15), box_w, Inches(0.75),
               big, size=36, bold=True, color=INK_900)
        _label(s, x, top + Inches(0.9), box_w, Inches(0.35),
               small, size=11, color=INK_500, font=MONO)

    # Red accent pill highlighting the demo-worthy number
    pill = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6) + (box_w + gap) * 0 + Inches(0.25), top + Inches(1.18),
        Inches(2.45), Inches(0.12),
    )
    try:
        pill.adjustments[0] = 0.5
    except Exception:
        pass
    pill.fill.solid()
    pill.fill.fore_color.rgb = RED
    pill.line.fill.background()
    pill.shadow.inherit = False

    # What ships
    _card(s, Inches(0.6), Inches(4.3), Inches(6.1), Inches(2.65))
    _textbox(s, Inches(0.9), Inches(4.4), Inches(5.7), Inches(0.4),
             "What ships, working end-to-end",
             size=15, bold=True, color=INK_900)
    _bullets(s, Inches(0.9), Inches(4.8), Inches(5.7), Inches(2.1), [
        "CLI:  cutpilot <source>  (local path or URL)",
        "HTTP:  cutpilot-serve  (FastAPI + uvicorn)",
        "UI:  ui/index.html  (static, Tailwind, Inter + JetBrains Mono)",
        "NAT path:  nat run --config_file=configs/cutpilot.yml",
        "Live 120 s CI smoke + 43 min opt-in e2e (`pytest -m e2e`)",
        "Lint:  ruff (py313)   ·   types:  mypy strict + pydantic plugin",
    ], size=12, line_spacing=1.2, bullet_color=INK_900)

    # What we cut (discipline signal)
    _card(s, Inches(6.85), Inches(4.3), Inches(5.9), Inches(2.65))
    _textbox(s, Inches(7.15), Inches(4.4), Inches(5.5), Inches(0.4),
             "What we cut on purpose",
             size=15, bold=True, color=INK_900)
    _bullets(s, Inches(7.15), Inches(4.8), Inches(5.5), Inches(2.1), [
        "VL-guided smart crop (kept 9:16 center crop).",
        "Scene detection tool.",
        "Audio normalization / fades.",
        "Word-level caption highlighting (full-segment only).",
        "Multi-speaker handling, Korean support.",
        "Separate Critic agent — Scout self-scores, Editor filters.",
    ], size=12, line_spacing=1.2, bullet_color=INK_300)

    _footer(s, 10, total)


def slide_thanks(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PAGE)

    _brand_dot(s, Inches(0.9), Inches(0.92), size=Inches(0.36))
    _textbox(s, Inches(1.45), Inches(0.8), Inches(10), Inches(0.6),
             "CutPilot", size=22, bold=True, color=INK_900)

    _textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2),
             "Thank you.", size=72, bold=True, color=INK_900)

    acc = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.7),
        Inches(1.4), Inches(0.08),
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = RED
    acc.line.fill.background()
    acc.shadow.inherit = False

    _card(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.95),
          fill=INK_900, border=INK_900)
    _textbox(s, Inches(1.2), Inches(4.25), Inches(11), Inches(0.4),
             "Try it", size=13, bold=True, color=INK_300, font=MONO)
    _textbox(s, Inches(1.2), Inches(4.6), Inches(11), Inches(0.5),
             "$  pip install -e “.[dev]”",
             size=20, color=YELLOW, font=MONO)
    _textbox(s, Inches(1.2), Inches(5.1), Inches(11), Inches(0.5),
             "$  cutpilot  sources/<your_video>.mp4",
             size=20, color=YELLOW, font=MONO)
    _textbox(s, Inches(1.2), Inches(5.6), Inches(11), Inches(0.5),
             "$  open  outputs/<run-id>/highlights.mp4",
             size=20, color=CARD, font=MONO)

    _textbox(s, Inches(0.9), Inches(6.5), Inches(12), Inches(0.5),
             "Questions?     sergey@lablup.com",
             size=16, color=INK_500, font=MONO)
    _footer(s, total, total)


def build() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_solution,
        slide_demo_flow,
        slide_architecture,
        slide_nemo_stack,
        slide_innovation,
        slide_impact,
        slide_demo_evidence,
        slide_completeness,
        slide_thanks,
    ]
    total = len(builders)
    for b in builders:
        b(prs, total)

    prs.save(str(OUT_PATH))
    return OUT_PATH


if __name__ == "__main__":
    path = build()
    print(f"wrote {path}")
